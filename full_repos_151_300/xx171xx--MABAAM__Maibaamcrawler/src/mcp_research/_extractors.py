"""Per-filetype extraction functions. Each returns (text, error_or_none)."""

import logging
import os
import re

logger = logging.getLogger(__name__)

# ── Text / Code ─────────────────────────────────────────────────────────────

_TEXT_EXTS = {
    ".txt", ".md", ".rst", ".csv", ".tsv", ".log", ".json", ".yaml", ".yml",
    ".xml", ".html", ".htm", ".css", ".js", ".jsx", ".ts", ".tsx", ".py",
    ".java", ".c", ".cpp", ".h", ".hpp", ".cs", ".go", ".rs", ".rb", ".php",
    ".sh", ".bat", ".ps1", ".lua", ".r", ".m", ".swift", ".kt", ".scala",
    ".sql", ".graphql", ".toml", ".ini", ".cfg", ".conf", ".env", ".nk",
    ".usda", ".usd", ".mtlx", ".mel", ".vex", ".osl", ".glsl", ".hlsl",
    ".vue", ".svelte", ".ex", ".exs", ".erl", ".hs", ".ml", ".clj", ".lisp",
    ".vim", ".el", ".tf", ".proto", ".dockerfile", ".makefile", ".cmake",
    ".gradle", ".sbt", ".pl", ".pm",
}

_PDF_EXTS = {".pdf"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
_AUDIO_EXTS = {".mp3", ".wav", ".ogg", ".flac", ".m4a", ".aac", ".opus", ".wma"}
_VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".webm", ".wmv", ".flv", ".m4v"}
_OFFICE_EXTS = {".docx", ".xlsx", ".pptx"}

ALL_SUPPORTED_EXTS = _TEXT_EXTS | _PDF_EXTS | _IMAGE_EXTS | _AUDIO_EXTS | _VIDEO_EXTS | _OFFICE_EXTS


def get_file_type(ext: str) -> str | None:
    """Classify extension into a type category. Returns None if unsupported."""
    ext = ext.lower()
    if ext in _TEXT_EXTS:
        return "text"
    if ext in _PDF_EXTS:
        return "pdf"
    if ext in _OFFICE_EXTS:
        return "office"
    if ext in _AUDIO_EXTS:
        return "audio"
    if ext in _VIDEO_EXTS:
        return "video"
    if ext in _IMAGE_EXTS:
        return "image"
    return None


def extract_text(path: str) -> tuple[str, str | None]:
    """Read text file with 500KB cap."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(500_000), None
    except Exception as e:
        return "", f"Read failed: {e}"


def extract_pdf(path: str) -> tuple[str, str | None]:
    """Extract text from PDF via PyPDF2."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "", "PyPDF2 not installed. Install: pip install 'mcp-research[ingest]'"
    try:
        reader = PdfReader(path)
        text = "\n\n".join(page.extract_text() or "" for page in reader.pages)
        return text, None
    except Exception as e:
        return "", f"PDF extraction failed: {e}"


def extract_docx(path: str) -> tuple[str, str | None]:
    """Extract text from DOCX."""
    try:
        from docx import Document
    except ImportError:
        return "", "python-docx not installed. Install: pip install 'mcp-research[ingest]'"
    try:
        doc = Document(path)
        text = "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
        return text, None
    except Exception as e:
        return "", f"DOCX extraction failed: {e}"


def extract_xlsx(path: str) -> tuple[str, str | None]:
    """Extract text from XLSX spreadsheets."""
    try:
        import openpyxl
    except ImportError:
        return "", "openpyxl not installed. Install: pip install 'mcp-research[ingest]'"
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts), None
    except Exception as e:
        return "", f"XLSX extraction failed: {e}"


def extract_pptx(path: str) -> tuple[str, str | None]:
    """Extract text from PPTX presentations."""
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx not installed. Install: pip install 'mcp-research[ingest]'"
    try:
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            if texts:
                parts.append(f"## Slide {i}\n" + "\n".join(texts))
        return "\n\n".join(parts), None
    except Exception as e:
        return "", f"PPTX extraction failed: {e}"


def extract_audio(path: str, whisper_model: str = "", whisper_device: str = "") -> tuple[str, str | None]:
    """Transcribe audio via faster-whisper."""
    try:
        from faster_whisper import WhisperModel
    except ImportError:
        return "", "faster-whisper not installed. Install: pip install 'mcp-research[youtube]'"
    from . import config
    w_model = whisper_model or config.WHISPER_MODEL or "base"
    w_device = whisper_device or config.WHISPER_DEVICE or "auto"
    try:
        device = w_device if w_device != "auto" else "cuda"
        try:
            wm = WhisperModel(w_model, device=device, compute_type="float16")
        except Exception:
            wm = WhisperModel(w_model, device="cpu", compute_type="int8")
        segs, _info = wm.transcribe(path, beam_size=5)
        text = " ".join(seg.text.strip() for seg in segs)
        return text, None
    except Exception as e:
        return "", f"Audio transcription failed: {e}"


def extract_video(path: str, whisper_model: str = "", whisper_device: str = "") -> tuple[str, str | None]:
    """Extract audio from video via ffmpeg, then transcribe."""
    import shutil
    import subprocess
    if not shutil.which("ffmpeg"):
        return "", "ffmpeg not found on PATH. Install ffmpeg for video transcription."
    tmp_audio = path + ".extract.wav"
    try:
        subprocess.run(
            ["ffmpeg", "-i", path, "-vn", "-acodec", "pcm_s16le", "-ar", "16000",
             "-ac", "1", "-y", tmp_audio],
            capture_output=True, timeout=300, check=True)
        return extract_audio(tmp_audio, whisper_model, whisper_device)
    except subprocess.SubprocessError as e:
        return "", f"ffmpeg failed: {e}"
    finally:
        try:
            os.remove(tmp_audio)
        except OSError:
            pass


def extract_image(path: str) -> tuple[str, str | None]:
    """Describe image via Ollama vision model."""
    from . import ollama
    text = ollama.ollama_describe_image(path)
    if text:
        return text, None
    from . import config
    if not config.OLLAMA_VISION_MODEL:
        return "", "No vision model configured. Set OLLAMA_VISION_MODEL env var."
    return "", "Image description failed."


def extract(path: str, whisper_model: str = "", whisper_device: str = "") -> tuple[str, str | None]:
    """Route to the appropriate extractor based on file extension."""
    ext = os.path.splitext(path)[1].lower()
    ftype = get_file_type(ext)
    if ftype is None:
        return "", f"Unsupported file type: {ext}"
    if ftype == "text":
        return extract_text(path)
    if ftype == "pdf":
        return extract_pdf(path)
    if ftype == "office":
        if ext == ".docx":
            return extract_docx(path)
        if ext == ".xlsx":
            return extract_xlsx(path)
        if ext == ".pptx":
            return extract_pptx(path)
        return "", f"Unsupported office format: {ext}"
    if ftype == "audio":
        return extract_audio(path, whisper_model, whisper_device)
    if ftype == "video":
        return extract_video(path, whisper_model, whisper_device)
    if ftype == "image":
        return extract_image(path)
    return "", f"Unknown type: {ftype}"
